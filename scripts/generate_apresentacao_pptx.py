"""Gera slides da apresentação executiva — Tech Challenge Fase 2."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "Apresentacao_Tech_Challenge_Fase2_Alfabetiza.pptx"

# Paleta institucional (educação / dados) — sem roxo genérico
NAVY = RGBColor(0x0B, 0x2C, 0x4A)
TEAL = RGBColor(0x0D, 0x7A, 0x6F)
ACCENT = RGBColor(0xE8, 0x7A, 0x2E)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
MUTED = RGBColor(0x5A, 0x6A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)


def _set_run(run, text: str, size: int, bold: bool = False, color=DARK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_bg(slide, color=LIGHT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _banner(slide, title: str, subtitle: str | None = None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    _set_run(p.add_run(), title, 28, True, WHITE)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.72), Inches(12.3), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        _set_run(sp.add_run(), subtitle, 14, False, RGBColor(0xB8, 0xC9, 0xD9))


def _bullets(slide, items: list[str], left=0.6, top=1.5, width=12, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        _set_run(p.add_run(), "•  " + item, size, False, DARK)


def _footer(slide, page: int, total: int):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    _set_run(
        p.add_run(),
        f"Tech Challenge Fase 2 · Alfabetiza-Cursor · nassereq/Alfabetiza-Cursor  ·  {page}/{total}",
        10,
        False,
        MUTED,
    )


def _card(slide, left, top, width, height, title: str, body: str, color=TEAL):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE0)
    tip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.12), Inches(height)
    )
    tip.fill.solid()
    tip.fill.fore_color.rgb = color
    tip.line.fill.background()
    t = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.15), Inches(width - 0.4), Inches(0.4))
    _set_run(t.text_frame.paragraphs[0].add_run(), title, 16, True, NAVY)
    b = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.55), Inches(width - 0.4), Inches(height - 0.7))
    bp = b.text_frame
    bp.word_wrap = True
    _set_run(bp.paragraphs[0].add_run(), body, 13, False, MUTED)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slides_meta: list = []

    # --- 1 capa ---
    s = prs.slides.add_slide(blank)
    _add_bg(s, NAVY)
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.8), Inches(13.333), Inches(1.7))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = TEAL
    stripe.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2))
    _set_run(t.text_frame.paragraphs[0].add_run(), "Pipeline Híbrida Medalhão", 40, True, WHITE)
    t2 = s.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.8))
    _set_run(
        t2.text_frame.paragraphs[0].add_run(),
        "Indicador Criança Alfabetizada — Tech Challenge Fase 2",
        24,
        False,
        RGBColor(0xC8, 0xE6, 0xE2),
    )
    t3 = s.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(11.5), Inches(1.0))
    tf = t3.text_frame
    _set_run(tf.paragraphs[0].add_run(), "POSTECH / FIAP — AI Scientist", 16, True, WHITE)
    p = tf.add_paragraph()
    _set_run(p.add_run(), "Engenharia de dados para política pública educacional · github.com/nassereq/Alfabetiza-Cursor", 14, False, WHITE)
    slides_meta.append(s)

    # --- 2 agenda ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "Agenda")
    _bullets(
        s,
        [
            "Introdução e objetivo da entrega",
            "Contextualização: Compromisso Nacional Criança Alfabetizada",
            "O problema de dados e o papel da engenharia",
            "Ferramentas e stack tecnológica",
            "O projeto: arquitetura híbrida e medalhão",
            "Resultados na camada Gold",
            "Qualidade, FinOps e evidência em nuvem",
            "Potencial de IA e otimizações futuras",
            "Repositório, entregáveis e fechamento",
        ],
        size=18,
    )
    slides_meta.append(s)

    # --- 3 introdução ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "Introdução", "Quem somos e o que entregamos")
    _card(
        s,
        0.5,
        1.6,
        6.0,
        2.2,
        "Papel da equipe",
        "Atuar como engenharia de dados de uma organização pública: construir a pipeline que transforma dados educacionais em evidência para gestão.",
        TEAL,
    )
    _card(
        s,
        6.8,
        1.6,
        6.0,
        2.2,
        "Objetivo do Tech Challenge",
        "Pipeline híbrida (batch + streaming) em arquitetura medalhão (Bronze → Silver → Gold), com qualidade, monitoramento, FinOps e preparação para IA.",
        ACCENT,
    )
    _card(
        s,
        0.5,
        4.1,
        6.0,
        2.2,
        "Entrega principal (90%)",
        "Apresentação executiva ≤ 5 min: problema, arquitetura, valor educacional e potencial de IA — linguagem para liderança, sem código na tela.",
        NAVY,
    )
    _card(
        s,
        6.8,
        4.1,
        6.0,
        2.2,
        "Artefato técnico (10%)",
        "Repositório GitHub público com evolução por branches/PRs, documentação, pipeline reproduzível e evidência cloud.",
        TEAL,
    )
    slides_meta.append(s)

    # --- 4 contexto ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "Contextualização", "Política pública de alfabetização")
    _bullets(
        s,
        [
            "Compromisso Nacional Criança Alfabetizada: União, estados, DF e municípios.",
            "Meta: toda criança alfabetizada até o final do 2º ano do ensino fundamental.",
            "Pesquisa Alfabetiza Brasil (INEP, 2023): ponto de corte 743 na escala Saeb.",
            "Indicador Criança Alfabetizada = % de estudantes no patamar de alfabetizados.",
            "Meta nacional: 100% até 2030 — exige monitoramento territorial contínuo.",
            "Fonte oficial usada: Base dos Dados / Inep (br_inep_avaliacao_alfabetizacao).",
        ],
    )
    slides_meta.append(s)

    # --- 5 problema ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "O problema", "Dados existem — decisão integrada não")
    _card(
        s,
        0.5,
        1.6,
        4.0,
        4.5,
        "Fragmentação",
        "Metas (Brasil/UF/município), território, indicador e microdados vivem em silos. Sem integração, a gestão não prioriza onde agir.",
        ACCENT,
    )
    _card(
        s,
        4.7,
        1.6,
        4.0,
        4.5,
        "Risco operacional",
        "Relatórios retrospectivos isolados não mostram gap meta × resultado nem evolução temporal comparável entre entes.",
        TEAL,
    )
    _card(
        s,
        8.9,
        1.6,
        4.0,
        4.5,
        "O que falta",
        "Uma pipeline confiável: rastreável (Bronze), tratada (Silver), analítica (Gold), com qualidade, custo controlado e frescor (streaming).",
        NAVY,
    )
    slides_meta.append(s)

    # --- 6 ferramentas ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "Ferramentas utilizadas", "Stack alinhada às aulas POSTECH")
    items = [
        ("Python / Pandas / PyArrow", "Pipeline local, transformação e Parquet"),
        ("Base dos Dados + BigQuery", "Origem oficial + evidência GCP (job_id)"),
        ("Arquitetura medalhão", "Bronze / Silver / Gold (SOR / SOT / SPEC)"),
        ("Kafka / file-sink", "Braço streaming de atualização do indicador"),
        ("MySQL + SQL", "Dimensões relacionais (aula de BD)"),
        ("AWS Glue template", "Caminho de deploy cloud da disciplina ETL"),
        ("GitHub + PRs", "Evolução visível (bronze → gold → cloud)"),
        ("Qualidade + FinOps", "validate.py, health check, estimativa de custo"),
    ]
    for i, (title, body) in enumerate(items):
        col = i % 4
        row = i // 4
        _card(s, 0.4 + col * 3.2, 1.5 + row * 2.5, 3.05, 2.2, title, body, TEAL if row == 0 else NAVY)
    slides_meta.append(s)

    # --- 7 arquitetura ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "Arquitetura da solução", "Pipeline híbrida batch + streaming")
    _bullets(
        s,
        [
            "Origens: Base dos Dados (BigQuery), dimensões MySQL, eventos JSON (estilo NoSQL).",
            "Batch: carga periódica de históricos (município/UF → entidades da pipeline).",
            "Streaming: producer de eventos do indicador (Kafka ou file-sink na demo).",
            "Bronze: bruto + metadados (_data_ingestao, _fonte) particionado por data.",
            "Silver: tipagem, dedup, FKs e integração indicador × meta × município.",
            "Gold: 3 produtos — indicador municipal, meta×resultado (UF/Brasil), evolução.",
            "Cloud: evidência real GCP BigQuery; template AWS S3 + Glue (aula ETL).",
        ],
        size=17,
    )
    slides_meta.append(s)

    # --- 8 projeto fluxo ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "O projeto em prática", "Do dado bruto à decisão")
    steps = [
        ("1. Ingestão", "fetch/map Base dos Dados\nserie=2 · rede=5\ntaxa → pct_alfabetizados"),
        ("2. Bronze", "Parquet bruto\nrastreabilidade\nbatch + eventos"),
        ("3. Silver", "Limpeza + chaves\nintegrado meta×ind\nqualidade OK"),
        ("4. Gold", "Município / UF / BR\nevolução temporal\nCSVs para BI"),
    ]
    for i, (title, body) in enumerate(steps):
        shape = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45 + i * 3.2), Inches(2.0), Inches(3.0), Inches(3.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = TEAL
        tb = s.shapes.add_textbox(Inches(0.6 + i * 3.2), Inches(2.2), Inches(2.7), Inches(0.5))
        _set_run(tb.text_frame.paragraphs[0].add_run(), title, 18, True, NAVY)
        bb = s.shapes.add_textbox(Inches(0.6 + i * 3.2), Inches(2.9), Inches(2.7), Inches(2.6))
        tf = bb.text_frame
        tf.word_wrap = True
        for j, line in enumerate(body.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            _set_run(p.add_run(), line, 14, False, MUTED)
        if i < 3:
            arr = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(3.35 + i * 3.2), Inches(3.6), Inches(0.35), Inches(0.35)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = ACCENT
            arr.line.fill.background()
    slides_meta.append(s)

    # --- 9 resultados Brasil ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "Resultados — Brasil", "Gold: comparativo meta × resultado")
    _card(s, 0.5, 1.6, 4.0, 4.8, "2023", "Taxa média: 60,47%\nMunicípios: 4.950\nNa meta: 2.959\nMeta traj.: 56%\nGap: +4,5 p.p.", TEAL)
    _card(s, 4.7, 1.6, 4.0, 4.8, "2024", "Taxa média: 63,17%\nMunicípios: 5.516\nNa meta: 2.968\nMeta traj.: 62%\nGap: +1,2 p.p.", ACCENT)
    _card(
        s,
        8.9,
        1.6,
        4.0,
        4.8,
        "Leitura",
        "Houve avanço da média nacional (60,5 → 63,2), mas o desafio 2030 (100%) permanece. A Gold torna o gap visível por território.",
        NAVY,
    )
    slides_meta.append(s)

    # --- 10 resultados UF ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "Resultados — desigualdade territorial", "2024 · taxa média por UF (amostra)")
    _card(s, 0.5, 1.6, 6.0, 2.4, "Maior desempenho (2024)", "Ceará 90,3% · Goiás 80,3% · Espírito Santo 78,6%\n(acima da meta didática de 62%)", TEAL)
    _card(s, 6.8, 1.6, 6.0, 2.4, "Maior desafio (2024)", "Bahia 36,6% · Sergipe 36,9% · Rio Grande do Norte 42,5%\n(gap de ~19 a 25 p.p. vs meta)", ACCENT)
    _card(
        s,
        0.5,
        4.3,
        12.3,
        2.2,
        "Valor para a gestão",
        "A mesma Gold responde: quem está abaixo da meta? qual a trajetória ano a ano? onde priorizar recomposição de aprendizagem? — de relatório estático para prioridade operacional.",
        NAVY,
    )
    slides_meta.append(s)

    # --- 11 produtos gold ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "Produtos analíticos (Gold)", "O que o enunciado pede — materializado")
    _bullets(
        s,
        [
            "indicador_alfabetizacao_municipio — fato municipal com meta, gap e flag atingiu_meta.",
            "comparativo_meta_resultado_uf / _brasil — agregações para gestão estadual e nacional.",
            "evolucao_temporal_* — delta em pontos percentuais vs ano anterior.",
            "Volume processado: ~10,5 mil linhas municipais (2023–2024, serie 2, rede 5).",
            "Saídas em Parquet + CSV (Excel / Power BI / notebooks).",
            "Qualidade automatizada: duplicidade, nulos, FKs e consistência meta×indicador.",
        ],
    )
    slides_meta.append(s)

    # --- 12 qualidade finops cloud ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "Qualidade, FinOps e nuvem", "Requisitos modernos do enunciado")
    _card(
        s,
        0.5,
        1.6,
        4.0,
        4.8,
        "Qualidade",
        "quality/validate.py\nSem duplicatas nas chaves\nFKs município→UF\nindicador→município\nAnos cobertos por meta\npassed = true",
        TEAL,
    )
    _card(
        s,
        4.7,
        1.6,
        4.0,
        4.8,
        "FinOps",
        "Parquet particionado\nCompute sob demanda\nDev local barato\nBQ ~0,7 MB/job\nAWS Glue sob demanda\n~US$ 5–20/mês (demo)",
        ACCENT,
    )
    _card(
        s,
        8.9,
        1.6,
        4.0,
        4.8,
        "Evidência cloud",
        "GCP BigQuery\nProjeto: alfabetiza-fiap-…\nJob ID versionado\nreports/cloud_evidence/\n+ template AWS Glue",
        NAVY,
    )
    slides_meta.append(s)

    # --- 13 IA ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "Uso de IA — potencial sobre a Gold", "Engenharia entrega o ativo; ciência consome")
    _bullets(
        s,
        [
            "Predição: projetar % alfabetizados por município até 2030 (features: histórico, UF, gap, rede).",
            "Desigualdade: decompor gaps regionais e correlacionar com vulnerabilidade socioeconômica (IBGE/CadÚnico — opcional).",
            "Clustering: grupos de municípios com perfil semelhante para políticas diferenciadas.",
            "Priorização: ranking de intervenção com base em distância da meta e tendência temporal.",
            "LLM / analytics assistida: perguntas em linguagem natural sobre a Gold (camada futura).",
            "Importante: esta entrega foca a plataforma de dados — modelos são o próximo ciclo.",
        ],
        size=17,
    )
    slides_meta.append(s)

    # --- 14 otimizações ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "Possíveis otimizações", "Roadmap técnico e de negócio")
    _bullets(
        s,
        [
            "Extrair microdados reais da tabela alunos (hoje: amostra sintética alinhada à taxa).",
            "Enriquecer nomes de município via diretório IBGE (Base dos Dados).",
            "Metas oficiais por ente (quando publicadas) no lugar da trajetória didática 2030.",
            "Orquestração (Airflow/Step Functions) e agendamento do batch.",
            "Deploy completo AWS: buckets SOR/SOT/SPEC + jobs Glue em produção.",
            "Streaming com Kafka gerenciado + schema registry; DLQ e replay.",
            "Camada semântica / dbt e catálogo (Data Hub / Unity Catalog).",
            "Observabilidade: alertas se quality.passed == false ou volume cair > 50%.",
        ],
        size=16,
    )
    slides_meta.append(s)

    # --- 15 entregáveis ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "Entregáveis e repositório", "O que enviar no AVA")
    _bullets(
        s,
        [
            "GitHub público: https://github.com/nassereq/Alfabetiza-Cursor",
            "PRs mergeados: Base dos Dados (#1) e evidência cloud GCP (#2).",
            "Pipeline: python -m pipelines.run_pipeline --fonte raw --with-streaming",
            "Docs: contexto, arquitetura, FinOps, dicionário, roteiro de vídeo e estágios.",
            "Evidência cloud: reports/cloud_evidence/ (job_id BigQuery).",
            "Este deck + vídeo executivo ≤ 5 min (problema → arquitetura → valor → IA).",
        ],
    )
    slides_meta.append(s)

    # --- 16 fechamento ---
    s = prs.slides.add_slide(blank)
    _add_bg(s)
    _banner(s, "Conclusão", "Evidência para a meta 2030")
    _bullets(
        s,
        [
            "Problema: alfabetização exige dados integrados para priorizar territórios.",
            "Solução: pipeline híbrida medalhão com qualidade, FinOps e cloud.",
            "Resultado: Gold pronta para gestão (município / UF / Brasil) e para IA.",
            "Próximo passo humano: gravar o vídeo com o roteiro em docs/ROTEIRO_VIDEO.md.",
            "Obrigado — perguntas?",
        ],
        size=20,
        top=1.8,
    )
    slides_meta.append(s)

    # --- 17 thanks ---
    s = prs.slides.add_slide(blank)
    _add_bg(s, NAVY)
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(1))
    p = t.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), "Obrigado", 48, True, WHITE)
    t2 = s.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.5), Inches(1.2))
    p2 = t2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    _set_run(p2.add_run(), "Alfabetiza-Cursor · Tech Challenge Fase 2 · FIAP POSTECH", 18, False, RGBColor(0xC8, 0xE6, 0xE2))
    t3 = s.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.6))
    p3 = t3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    _set_run(p3.add_run(), "github.com/nassereq/Alfabetiza-Cursor", 16, False, WHITE)
    slides_meta.append(s)

    total = len(slides_meta)
    for i, slide in enumerate(slides_meta, start=1):
        if i == 1 or i == total:
            continue
        _footer(slide, i, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Salvo: {path}")
